import os
import requests
import uuid
import shutil
import pytesseract  
import csv

from django.conf import settings
from sentry_sdk import capture_exception, capture_message
from django.utils import timezone
from django.core.paginator import Paginator

from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
from rest_framework.permissions import IsAuthenticated
from .authentication import CustomTokenAuthentication

from PyPDF2 import PdfReader
from pdf2image import convert_from_path
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
from datetime import datetime
from contextlib import contextmanager



from simple_salesforce import Salesforce

from users.models import SalesforceConnection, APIUsage, APIKey, ProcessLog
from .permissions import IsSuperUser

pytesseract.tesseract_cmd = "/usr/bin/tesseract"


class SuperUserOnlyView(APIView):
    authentication_classes = [CustomTokenAuthentication]
    permission_classes = [IsSuperUser]

    def get(self, request, *args, **kwargs):
        data = {"message": "This is a superuser-only API"}
        return Response(data, status=status.HTTP_200_OK)

class UserAPIUsageByTransactionView(APIView):
    authentication_classes = [CustomTokenAuthentication]
    permission_classes = [IsAuthenticated]

    def get(self, request, *args, **kwargs):
        try:
            transaction_id = request.query_params.get('transaction_id')

            if not transaction_id:
                return Response({"detail": "transaction_id is required."}, status=status.HTTP_400_BAD_REQUEST)

            api_usage = APIUsage.objects.filter(transaction_id=transaction_id).first()

            if not api_usage:
                return Response({"detail": "APIUsage record not found for the provided transaction_id."}, status=status.HTTP_404_NOT_FOUND)

            process_logs = ProcessLog.objects.filter(api_usage=api_usage)

            data = self.get_api_usage_serializer(api_usage, process_logs)

            return Response(data, status=status.HTTP_200_OK)

        except Exception as e:
            sentry_sdk.capture_exception(e)
            return Response({"detail": "An error occurred while fetching the logs."}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    def get_api_usage_serializer(self, api_usage, process_logs):
        process_logs_data = [
            {
                "step_name": process_log.step_name,
                "start_time": process_log.start_time,
                "end_time": process_log.end_time,
                "duration_seconds": process_log.duration_seconds,
                "status": process_log.status,
                "error_message": process_log.error_message,
            }
            for process_log in process_logs
        ]

        return {
            "status": api_usage.process_status,
            "timestamp": api_usage.timestamp,
            "transaction_id": api_usage.transaction_id,
            "error_message": api_usage.error_message,
            "process_duration": api_usage.process_duration,
            "document_id": api_usage.sf_document_id,
            "process_step_logs": process_logs_data,
        }

class UserAPIUsageLogsView(APIView):
    authentication_classes = [CustomTokenAuthentication]
    permission_classes = [IsAuthenticated]

    def get(self, request, *args, **kwargs):
        try:
            user = request.user

            page_size = int(request.query_params.get('page_size', 10))
            page_number = int(request.query_params.get('page', 1))

            api_usage_logs = APIUsage.objects.filter(user=user).order_by('-timestamp')

            paginator = Paginator(api_usage_logs, page_size)
            paginated_logs = paginator.get_page(page_number)

            data = self.get_api_usage_logs_serializer(paginated_logs)

            return Response({
                'total_record_count': paginator.count,
                'api_logs': data
            }, status=status.HTTP_200_OK)

        except Exception as e:
            # Log the exception in Sentry
            sentry_sdk.capture_exception(e)
            return Response({"detail": "An error occurred while fetching the logs."}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    def get_api_usage_logs_serializer(self, logs):
        result = []
        for log in logs:
            process_logs = ProcessLog.objects.filter(api_usage=log)
            process_logs_data = [
                {
                    "step_name": process_log.step_name,
                    "start_time": process_log.start_time,
                    "end_time": process_log.end_time,
                    "duration_seconds": process_log.duration_seconds,
                    "status": process_log.status,
                    "error_message": process_log.error_message,
                }
                for process_log in process_logs
            ]

            result.append({
                "status": log.process_status,
                "timestamp": log.timestamp,
                "transaction_id": log.transaction_id,
                "error_message": log.error_message,
                "process_duration": log.process_duration,
                "document_id": log.sf_document_id,
                "process_step_logs": process_logs_data,
            })

        return result

class DocumentProcessingView(APIView):
    authentication_classes = [CustomTokenAuthentication]
    permission_classes = [IsAuthenticated]

    def post(self, request):
        document_id = request.data.get("documentId")
        organisation_id = request.data.get("organisationId")
        user_id = str(request.user.id)  # Get user ID
        transaction_id = str(uuid.uuid4())

        if not document_id or not organisation_id:
            self.log_api_usage(request.user, None, document_id, "FAILURE", request, transaction_id, 'Missing documentId or organisationId')
            return Response(
                {"error": "Missing documentId or organisationId"},
                status=status.HTTP_400_BAD_REQUEST,
            )

        api_usage = self.log_api_usage(request.user, None, document_id, "PROCESSING", request, transaction_id, '')
        api_usage.process_start_time = timezone.now()
        api_usage.process_status = "PROCESSING"
        api_usage.save()

        file_path = None
        transaction_dir = None
        try:
            user_dir = os.path.join(settings.MEDIA_ROOT, str(user_id))
            if not os.path.exists(user_dir):
                os.makedirs(user_dir)

            transaction_dir = os.path.join(user_dir, transaction_id)
            if not os.path.exists(transaction_dir):
                os.makedirs(transaction_dir)

            with self.process_step(api_usage, 'Fetch Salesforce Connection'):
                connection = SalesforceConnection.objects.get(
                    user=request.user,
                    organization_id=organisation_id
                )

            api_usage.salesforce_connection = connection
            api_usage.save()

            with self.process_step(api_usage, 'Fetch File from Salesforce'):
                file_path = self.fetch_file_from_salesforce(connection.access_token, document_id, connection.instance_url, transaction_dir)

            with self.process_step(api_usage, 'Process File'):
                parsed_text, num_pages, num_characters = self.process_file(file_path)
            
            response_data = {
                "transactionId": transaction_id,
                "fileName": os.path.basename(file_path),
                "numPages": num_pages,
                "numCharacters": num_characters,
                "parsedText": parsed_text,
            }
            
            api_usage.process_end_time = timezone.now()
            api_usage.process_status = "SUCCESS"
            api_usage.calculate_process_duration()
            api_usage.save()

            return Response(response_data, status=status.HTTP_200_OK)
        except Exception as e:
            capture_exception(e)
            api_usage.process_end_time = timezone.now()
            api_usage.process_status = "FAILURE"
            api_usage.calculate_process_duration()
            api_usage.save()
            return Response({"error": str(e), "transactionId": transaction_id}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
        finally:
            if transaction_dir and os.path.exists(transaction_dir):
                shutil.rmtree(transaction_dir)

    def fetch_file_from_salesforce(self, access_token, document_id, instance_url, transaction_dir):
        sf = Salesforce(instance_url=instance_url, session_id=access_token)
        query = f"SELECT VersionData, Title, FileExtension FROM ContentVersion WHERE Id = '{document_id}'"
        content_version = sf.query(query)
        if not content_version["records"]:
            raise ValueError(f"No file found for DocumentId {document_id}")
        
        version_data_relative_url = content_version["records"][0]["VersionData"]
        file_name = content_version["records"][0]["Title"]
        file_extension = content_version["records"][0]["FileExtension"]
        
        version_data_url = f"{instance_url}{version_data_relative_url}"
        headers = {"Authorization": f"Bearer {access_token}"}
        response = requests.get(version_data_url, headers=headers, stream=True)
        if response.status_code != 200:
            raise ValueError(f"Failed to fetch file content. HTTP Status {response.status_code}")
        
        file_path = os.path.join(transaction_dir, f"{file_name}.{file_extension}")
        with open(file_path, "wb") as f:
            for chunk in response.iter_content(chunk_size=1024):
                f.write(chunk)
        return file_path

    def process_file(self, file_path):
        file_extension = os.path.splitext(file_path)[1].lower()
        parsed_text = None
        num_pages = 0
        num_characters = 0
        if file_extension == ".pdf":
            parsed_text, num_pages, num_characters = self.extract_text_from_pdf(file_path)
        elif file_extension == ".docx":
            parsed_text = self.extract_text_from_docx(file_path)
        elif file_extension == ".xlsx":
            parsed_text = self.extract_text_from_xlsx(file_path)
        elif file_extension == ".pptx":
            parsed_text = self.extract_text_from_pptx(file_path)
        elif file_extension == ".csv":
            parsed_text = self.extract_text_from_csv(file_path)
        num_characters = len(parsed_text or "")
        return parsed_text, num_pages, num_characters

    def extract_text_from_pdf(self, file_path):
        reader = PdfReader(file_path)
        text = ""
        num_pages = len(reader.pages)
        for page_number, page in enumerate(reader.pages):
            extracted_text = page.extract_text()

            if extracted_text and extracted_text.strip():
                text += extracted_text
            else:
                if num_pages > 3:
                    capture_message('more than 3 pages being run at dpi 100')
                    text += self.ocr_pdf_page(file_path, page_number, 100)
                else:
                    capture_message('less than 3 pages being run at dpi 300')
                    text += self.ocr_pdf_page(file_path, page_number, 300)
        num_characters = len(text)
        return text, num_pages, num_characters

    def ocr_pdf_page(self, file_path, page_number, dpi):
        try:
            images = convert_from_path(file_path, first_page=page_number + 1, last_page=page_number + 1, dpi=dpi)
            if not images:
                message = f"No images generated for page {page_number + 1} from PDF."
                capture_exception(Exception(message))
                return ""

            ocr_text = ""
            for image_index, image in enumerate(images):
                try:
                    if image.mode != 'RGB':
                        image = image.convert('RGB')

                    image_path = os.path.join(settings.MEDIA_ROOT, f"ocr_image_page_{page_number + 1}_image_{image_index + 1}.png")
                    image.save(image_path)
                    ocr_data = pytesseract.image_to_data(image, config='--psm 6', nice=1, output_type=pytesseract.Output.DICT, timeout=180)
                    ocr_text += ' '.join([ocr_data['text'][i] for i in range(len(ocr_data['text'])) if int(ocr_data['conf'][i]) > 0])
                except Exception as inner_e:
                    message = f"Image handling failed on page {page_number + 1}, image {image_index + 1}: {str(inner_e)}"
                    capture_exception(Exception(message))
                    ocr_text += f"[Image Handling Error: {message}]\n"

            return ocr_text
        except Exception as e:
            message = f"OCR process failed for page {page_number + 1}: {str(e)}"
            capture_exception(Exception(message))
            return f"[OCR Process Error: {message}]"

    def extract_text_from_docx(self, file_path):
        document = Document(file_path)
        text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        return text

    def extract_text_from_xlsx(self, file_path):
        workbook = load_workbook(file_path)
        text = ""
        for sheet in workbook:
            for row in sheet.iter_rows(values_only=True):
                text += " ".join(str(cell) for cell in row if cell) + "\n"
        return text

    def extract_text_from_pptx(self, file_path):
        presentation = Presentation(file_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def extract_text_from_csv(self, file_path):
        text = ""
        with open(file_path, newline='', encoding="utf-8") as csvfile:
            reader = csv.reader(csvfile)
            for row in reader:
                text += " ".join(row) + "\n"
        return text

    def log_api_usage(self, user, connection, document_id, status, request, transaction_id, message):
        try:
            token = self.get_token_from_request(request)
            api_key = APIKey.objects.filter(key=token).first()

            api_usage = APIUsage.objects.create(
                user=user,
                api_key=api_key,
                salesforce_connection=connection,
                sf_document_id=document_id,
                status=status,
                process_status=status,
                transaction_id=transaction_id,
                error_message=message,
            )
            return api_usage
        except Exception as e:
            capture_exception(e)

    def get_token_from_request(self, request):
        auth_header = request.headers.get("Authorization")
        if auth_header and auth_header.startswith("Token "):
            return auth_header.split()[1]
        return None

    @contextmanager
    def process_step(self, api_usage, step_name):
        process_log = ProcessLog.objects.create(
            api_usage=api_usage,
            step_name=step_name,
            start_time=timezone.now(),
            status='PROCESSING'
        )
        try:
            yield
            process_log.end_time = timezone.now()
            process_log.status = 'SUCCESS'
            process_log.calculate_duration()
            process_log.save()
        except Exception as e:
            process_log.end_time = timezone.now()
            process_log.status = 'FAILURE'
            process_log.error_message = str(e)
            process_log.calculate_duration()
            process_log.save()
            raise
